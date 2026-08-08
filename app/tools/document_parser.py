"""Extract plain text from uploaded .pdf / .docx / .pptx documents.

The downstream pipeline (content parser, plan generator, etc.) only consumes
plain text, so this tool is the single "file -> text" adapter at the front of
the agent workflow. It never writes the original file to disk: callers only
store the sanitized basename and the extracted text.
"""

import io
import re
from dataclasses import dataclass

SUPPORTED_EXTENSIONS = (".pdf", ".docx", ".pptx")


@dataclass(frozen=True)
class DocumentParseResult:
    text: str
    file_format: str  # "pdf" | "docx" | "pptx"
    char_count: int


class DocumentParseError(Exception):
    """User-facing error raised when an upload cannot be turned into text."""


def detect_format(filename: str, data: bytes) -> str | None:
    """Detect the document format from the extension whitelist and magic bytes.

    Returns ``"pdf"``, ``"docx"`` or ``"pptx"``, or ``None`` when the file is
    not a supported, well-formed document. The extension alone is never
    trusted: PDFs must start with ``%PDF-`` and DOCX/PPTX are ZIP containers
    starting with ``PK``.
    """
    extension = _extension(filename)
    if extension == ".pdf":
        return "pdf" if data[:5] == b"%PDF-" else None
    if extension in (".docx", ".pptx"):
        return extension[1:] if data[:2] == b"PK" else None
    return None


def extract_text_from_bytes(filename: str, data: bytes) -> DocumentParseResult:
    """Extract sanitized plain text from the given document bytes."""
    if not data:
        raise DocumentParseError("文件内容为空。")

    file_format = detect_format(filename, data)
    if file_format is None:
        raise DocumentParseError("不支持的文件格式，仅支持 .pdf / .docx / .pptx 文本型文件。")

    try:
        if file_format == "pdf":
            text = _extract_pdf(data)
        elif file_format == "docx":
            text = _extract_docx(data)
        else:
            text = _extract_pptx(data)
    except DocumentParseError:
        raise
    except Exception as exc:  # corrupt container, encrypted pdf, etc.
        raise DocumentParseError("无法读取文件内容，请确认文件未损坏或未加密。") from exc

    text = _sanitize_text(text)
    return DocumentParseResult(text=text, file_format=file_format, char_count=len(text))


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(parts)


def _extract_docx(data: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append("、".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append("、".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _sanitize_text(text: str) -> str:
    """Drop control characters, collapse spaces, and remove blank lines."""
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line)


def _extension(filename: str) -> str:
    name = filename.rsplit("/", 1)[-1].rsplit("\\", 1)[-1]
    dot = name.rfind(".")
    if dot <= 0:
        return ""
    return name[dot:].lower()
