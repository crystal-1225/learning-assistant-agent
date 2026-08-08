import io

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas

from app.tools.document_parser import (
    DocumentParseError,
    detect_format,
    extract_text_from_bytes,
)


def make_docx(lines: list[str]) -> bytes:
    buf = io.BytesIO()
    doc = Document()
    for line in lines:
        doc.add_paragraph(line)
    doc.save(buf)
    return buf.getvalue()


def make_pptx(lines: list[str]) -> bytes:
    buf = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
    box.text_frame.text = lines[0]
    for line in lines[1:]:
        paragraph = box.text_frame.add_paragraph()
        paragraph.text = line
    prs.save(buf)
    return buf.getvalue()


def make_pdf(lines: list[str]) -> bytes:
    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.setFont("Helvetica", 12)
    y = 760
    for line in lines:
        c.drawString(72, y, line)
        y -= 20
    c.save()
    return buf.getvalue()


def test_detect_format_by_magic_bytes() -> None:
    assert detect_format("notes.pdf", b"%PDF-1.7 example") == "pdf"
    assert detect_format("a.docx", b"PK\x03\x04 zip-content") == "docx"
    assert detect_format("a.pptx", b"PK\x03\x04 zip-content") == "pptx"
    assert detect_format("renamed.pdf", b"PK\x03\x04 zip-content") is None
    assert detect_format("notes.txt", b"plain text") is None
    assert detect_format("", b"") is None


def test_extract_docx_text() -> None:
    result = extract_text_from_bytes("notes.docx", make_docx(["顺序表", "单链表", " 栈 "]))
    assert result.file_format == "docx"
    assert "顺序表" in result.text
    assert "单链表" in result.text
    assert "栈" in result.text
    assert result.char_count == len(result.text)


def test_extract_pptx_text() -> None:
    result = extract_text_from_bytes("slides.pptx", make_pptx(["数据结构", "顺序表、栈、队列"]))
    assert result.file_format == "pptx"
    assert "数据结构" in result.text
    assert "顺序表" in result.text
    assert "队列" in result.text


def test_extract_pdf_text() -> None:
    result = extract_text_from_bytes("notes.pdf", make_pdf(["Data Structures", "Linked List"]))
    assert result.file_format == "pdf"
    assert "Data Structures" in result.text


def test_uppercase_extension_is_accepted() -> None:
    result = extract_text_from_bytes("NOTES.DOCX", make_docx(["内容"]))
    assert result.file_format == "docx"


def test_unsupported_extension_is_rejected() -> None:
    with pytest.raises(DocumentParseError, match="不支持的文件格式"):
        extract_text_from_bytes("notes.txt", b"plain text")


def test_renamed_file_is_rejected_by_magic_bytes() -> None:
    with pytest.raises(DocumentParseError, match="不支持的文件格式"):
        extract_text_from_bytes("notes.pdf", b"PK\x03\x04 not actually a pdf")


def test_empty_file_is_rejected() -> None:
    with pytest.raises(DocumentParseError, match="文件内容为空"):
        extract_text_from_bytes("notes.docx", b"")


def test_corrupt_container_is_rejected() -> None:
    with pytest.raises(DocumentParseError, match="无法读取文件内容"):
        extract_text_from_bytes("notes.docx", b"PK\x03\x04 corrupt-not-a-real-docx")
